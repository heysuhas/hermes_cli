"""Structured, bounded model tools for local PDF and Office workflows."""

from __future__ import annotations

import hashlib
import json
import os
import shutil
import time
import uuid
from pathlib import Path
from typing import Any

from agent.corporate_events import (
    ComplianceEvent,
    OperationalEvent,
    record_compliance,
    record_operational,
)
from agent.corporate_policy import get_corporate_policy
from hermes_constants import get_hermes_home
from tools.registry import tool_error, tool_result


SUPPORTED_SUFFIXES = {".pdf", ".docx", ".xlsx", ".xlsm", ".pptx"}
PLAN_DIR = get_hermes_home() / "corporate" / "office-plans"


DOCUMENT_INSPECT_SCHEMA = {
    "name": "document_inspect",
    "description": (
        "Inspect a local PDF, Word, Excel, or PowerPoint file and return bounded "
        "metadata and structure. Use before extraction or planning edits."
    ),
    "parameters": {
        "type": "object",
        "properties": {"path": {"type": "string"}},
        "required": ["path"],
        "additionalProperties": False,
    },
}

DOCUMENT_EXTRACT_SCHEMA = {
    "name": "document_extract",
    "description": (
        "Extract bounded local content from PDF, DOCX, XLSX/XLSM, or PPTX. "
        "Supports page, sheet, and slide selection without using the internet."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "path": {"type": "string"},
            "selector": {
                "type": "string",
                "description": (
                    "Optional page range such as '1-5', Excel sheet name, or "
                    "PowerPoint slide range such as '1,3-5'."
                ),
            },
            "max_chars": {
                "type": "integer",
                "minimum": 1000,
                "maximum": 50000,
                "default": 12000,
            },
        },
        "required": ["path"],
        "additionalProperties": False,
    },
}

OFFICE_PLAN_SCHEMA = {
    "name": "office_plan_changes",
    "description": (
        "Validate and stage changes to an existing DOCX, XLSX/XLSM, or PPTX. "
        "Returns a structured preview and one-time plan_id; it does not modify "
        "the file. Word operations: replace_text, append_paragraph. Excel: "
        "set_cell, set_formula. PowerPoint: replace_text, add_text_slide."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "path": {"type": "string"},
            "output_path": {
                "type": "string",
                "description": (
                    "Optional approved output path. Omit to modify the original "
                    "after approval while creating a backup."
                ),
            },
            "operations": {
                "type": "array",
                "minItems": 1,
                "maxItems": 50,
                "items": {
                    "type": "object",
                    "properties": {
                        "action": {"type": "string"},
                        "find": {"type": "string"},
                        "replace": {"type": "string"},
                        "text": {"type": "string"},
                        "sheet": {"type": "string"},
                        "cell": {"type": "string"},
                        "value": {},
                        "formula": {"type": "string"},
                        "title": {"type": "string"},
                        "body": {"type": "string"},
                    },
                    "required": ["action"],
                    "additionalProperties": False,
                },
            },
        },
        "required": ["path", "operations"],
        "additionalProperties": False,
    },
}

OFFICE_APPLY_SCHEMA = {
    "name": "office_apply_changes",
    "description": (
        "Request user approval, then apply one previously staged Office change "
        "plan. Creates a backup for in-place edits, verifies the output, and "
        "invalidates the plan_id after successful use."
    ),
    "parameters": {
        "type": "object",
        "properties": {"plan_id": {"type": "string"}},
        "required": ["plan_id"],
        "additionalProperties": False,
    },
}


def _resolve_document(path_value: str, *, must_exist: bool = True) -> Path:
    policy = get_corporate_policy()
    path = Path(path_value).expanduser().resolve(strict=False)
    from agent.corporate_path_access import request_path_access

    error = request_path_access(
        path,
        purpose=(
            "read and process this document"
            if must_exist
            else "create or update a document in this folder"
        ),
    )
    if error:
        raise PermissionError(error)
    if path.suffix.lower() not in SUPPORTED_SUFFIXES:
        raise ValueError(
            "Supported formats are PDF, DOCX, XLSX/XLSM, and PPTX."
        )
    if must_exist and not path.is_file():
        raise FileNotFoundError(f"Document not found: {path}")
    if must_exist and path.stat().st_size > policy.max_document_bytes:
        raise ValueError(
            f"Document exceeds corporate size limit ({policy.max_document_bytes} bytes)."
        )
    return path


def _parse_numbers(selector: str | None, maximum: int) -> list[int]:
    if not selector:
        return list(range(1, maximum + 1))
    values: set[int] = set()
    for part in selector.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            left, right = part.split("-", 1)
            start, end = int(left), int(right)
            values.update(range(start, end + 1))
        else:
            values.add(int(part))
    return [value for value in sorted(values) if 1 <= value <= maximum]


def _inspect(path: Path) -> dict[str, Any]:
    suffix = path.suffix.lower()
    result: dict[str, Any] = {
        "path": str(path),
        "format": suffix.lstrip("."),
        "size_bytes": path.stat().st_size,
        "modified_at": path.stat().st_mtime,
    }
    if suffix == ".pdf":
        import fitz

        with fitz.open(path) as document:
            result.update(
                pages=document.page_count,
                metadata={key: value for key, value in document.metadata.items() if value},
                encrypted=document.is_encrypted,
            )
    elif suffix == ".docx":
        from docx import Document

        document = Document(path)
        result.update(
            paragraphs=len(document.paragraphs),
            tables=len(document.tables),
            headings=[
                paragraph.text[:200]
                for paragraph in document.paragraphs
                if paragraph.style and paragraph.style.name.startswith("Heading")
            ][:50],
        )
    elif suffix in {".xlsx", ".xlsm"}:
        from openpyxl import load_workbook

        workbook = load_workbook(path, read_only=True, data_only=False)
        try:
            result["sheets"] = [
                {
                    "name": sheet.title,
                    "rows": sheet.max_row,
                    "columns": sheet.max_column,
                }
                for sheet in workbook.worksheets
            ]
        finally:
            workbook.close()
    else:
        from pptx import Presentation

        presentation = Presentation(path)
        result.update(
            slides=len(presentation.slides),
            slide_titles=[
                (
                    slide.shapes.title.text[:200]
                    if slide.shapes.title is not None
                    else ""
                )
                for slide in presentation.slides
            ],
        )
    return result


def _extract(path: Path, selector: str | None, max_chars: int) -> dict[str, Any]:
    suffix = path.suffix.lower()
    chunks: list[dict[str, Any]] = []
    used = 0

    def add_chunk(label: str, text: str) -> bool:
        nonlocal used
        remaining = max_chars - used
        if remaining <= 0:
            return False
        bounded = text[:remaining]
        chunks.append({"source": label, "text": bounded})
        used += len(bounded)
        return len(text) <= remaining

    if suffix == ".pdf":
        import fitz

        with fitz.open(path) as document:
            for page_number in _parse_numbers(selector, document.page_count):
                text = document.load_page(page_number - 1).get_text("text")
                if not text.strip():
                    pixmap = document.load_page(page_number - 1).get_pixmap(
                        matrix=fitz.Matrix(2, 2),
                        alpha=False,
                    )
                    text = _windows_ocr_png(pixmap.tobytes("png"))
                    if not text.strip():
                        text = (
                            "[No embedded text and Windows local OCR was "
                            "unavailable or returned no text.]"
                        )
                if not add_chunk(f"page:{page_number}", text):
                    break
    elif suffix == ".docx":
        from docx import Document

        document = Document(path)
        for index, paragraph in enumerate(document.paragraphs, start=1):
            if paragraph.text.strip() and not add_chunk(
                f"paragraph:{index}", paragraph.text
            ):
                break
        for table_index, table in enumerate(document.tables, start=1):
            rows = [
                "\t".join(cell.text for cell in row.cells)
                for row in table.rows
            ]
            if rows and not add_chunk(f"table:{table_index}", "\n".join(rows)):
                break
    elif suffix in {".xlsx", ".xlsm"}:
        from openpyxl import load_workbook

        workbook = load_workbook(path, read_only=True, data_only=False)
        try:
            names = [selector] if selector else workbook.sheetnames
            for name in names:
                if name not in workbook.sheetnames:
                    raise ValueError(f"Worksheet not found: {name}")
                sheet = workbook[name]
                lines: list[str] = []
                for row in sheet.iter_rows(
                    min_row=1,
                    max_row=min(sheet.max_row, 500),
                    values_only=True,
                ):
                    lines.append("\t".join("" if value is None else str(value) for value in row))
                if not add_chunk(f"sheet:{name}", "\n".join(lines)):
                    break
        finally:
            workbook.close()
    else:
        from pptx import Presentation

        presentation = Presentation(path)
        for slide_number in _parse_numbers(selector, len(presentation.slides)):
            slide = presentation.slides[slide_number - 1]
            text = "\n".join(
                shape.text
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            )
            if not add_chunk(f"slide:{slide_number}", text):
                break
    return {
        "path": str(path),
        "format": suffix.lstrip("."),
        "chunks": chunks,
        "characters": used,
        "truncated": used >= max_chars,
    }


def _windows_ocr_png(image_bytes: bytes) -> str:
    """Run the packaged Windows.Media.Ocr bindings without network access."""
    if os.name != "nt":
        return ""
    try:
        import asyncio
        from winrt.windows.graphics.imaging import BitmapDecoder
        from winrt.windows.media.ocr import OcrEngine
        from winrt.windows.storage.streams import DataWriter, InMemoryRandomAccessStream
    except ImportError:
        return ""

    async def recognize() -> str:
        stream = InMemoryRandomAccessStream()
        writer = DataWriter(stream)
        writer.write_bytes(image_bytes)
        await writer.store_async()
        await writer.flush_async()
        writer.detach_stream()
        stream.seek(0)
        decoder = await BitmapDecoder.create_async(stream)
        bitmap = await decoder.get_software_bitmap_async()
        engine = OcrEngine.try_create_from_user_profile_languages()
        if engine is None:
            return ""
        result = await engine.recognize_async(bitmap)
        return str(result.text or "")

    try:
        return asyncio.run(recognize())
    except Exception:
        return ""


def _validate_operations(path: Path, operations: list[dict[str, Any]]) -> None:
    allowed = {
        ".docx": {"replace_text", "append_paragraph"},
        ".xlsx": {"set_cell", "set_formula"},
        ".xlsm": {"set_cell", "set_formula"},
        ".pptx": {"replace_text", "add_text_slide"},
    }
    if path.suffix.lower() not in allowed:
        raise ValueError("PDF files are read-only in corporate V1.")
    for operation in operations:
        action = str(operation.get("action") or "")
        if action not in allowed[path.suffix.lower()]:
            raise ValueError(
                f"Unsupported {path.suffix.lower()} operation: {action}"
            )
        if action == "replace_text" and not str(operation.get("find") or ""):
            raise ValueError("replace_text requires a non-empty find value.")
        if action in {"set_cell", "set_formula"}:
            if not operation.get("sheet") or not operation.get("cell"):
                raise ValueError(f"{action} requires sheet and cell.")


def _plan_path(plan_id: str) -> Path:
    if not plan_id or not plan_id.replace("-", "").isalnum():
        raise ValueError("Invalid plan_id.")
    return PLAN_DIR / f"{plan_id}.json"


def _plan_changes(path: Path, output_path: Path | None, operations: list[dict[str, Any]]) -> dict[str, Any]:
    _validate_operations(path, operations)
    plan_id = uuid.uuid4().hex
    source_fingerprint = hashlib.sha256(path.read_bytes()).hexdigest()
    record = {
        "version": 1,
        "plan_id": plan_id,
        "path": str(path),
        "output_path": str(output_path) if output_path else "",
        "operations": operations,
        "source_sha256": source_fingerprint,
        "created_at": time.time(),
        "expires_at": time.time() + 3600,
        "used": False,
    }
    PLAN_DIR.mkdir(parents=True, exist_ok=True)
    target = _plan_path(plan_id)
    temporary = target.with_suffix(".tmp")
    temporary.write_text(json.dumps(record, indent=2, ensure_ascii=False), encoding="utf-8")
    os.replace(temporary, target)
    return {
        "plan_id": plan_id,
        "expires_at": record["expires_at"],
        "source": str(path),
        "output": str(output_path or path),
        "format": path.suffix.lstrip("."),
        "operations": operations,
        "requires_approval": True,
    }


def _ask_approval(record: dict[str, Any]) -> bool:
    try:
        from tools.approval import prompt_dangerous_approval
        from tools.terminal_tool import _get_approval_callback

        choice = prompt_dangerous_approval(
            f"Apply Office plan {record['plan_id']} to {record['path']}",
            (
                f"Modify an existing Office document with "
                f"{len(record['operations'])} planned operation(s)"
            ),
            allow_permanent=False,
            approval_callback=_get_approval_callback(),
        )
        return choice in {"once", "session"}
    except Exception:
        return False


def _apply_operations(target: Path, operations: list[dict[str, Any]]) -> str:
    suffix = target.suffix.lower()
    if os.name == "nt":
        try:
            return _apply_with_com(target, operations)
        except Exception:
            pass
    if suffix == ".docx":
        from docx import Document

        document = Document(target)
        for operation in operations:
            if operation["action"] == "replace_text":
                for paragraph in document.paragraphs:
                    if operation["find"] in paragraph.text:
                        paragraph.text = paragraph.text.replace(
                            operation["find"], str(operation.get("replace") or "")
                        )
            else:
                document.add_paragraph(str(operation.get("text") or ""))
        document.save(target)
        return "python-docx"
    if suffix in {".xlsx", ".xlsm"}:
        from openpyxl import load_workbook

        workbook = load_workbook(target, keep_vba=suffix == ".xlsm")
        for operation in operations:
            sheet = workbook[str(operation["sheet"])]
            cell = sheet[str(operation["cell"])]
            cell.value = (
                operation.get("formula")
                if operation["action"] == "set_formula"
                else operation.get("value")
            )
        workbook.save(target)
        workbook.close()
        return "openpyxl"
    if suffix == ".pptx":
        from pptx import Presentation

        presentation = Presentation(target)
        for operation in operations:
            if operation["action"] == "replace_text":
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and operation["find"] in shape.text:
                            shape.text = shape.text.replace(
                                operation["find"], str(operation.get("replace") or "")
                            )
            else:
                layout = presentation.slide_layouts[1]
                slide = presentation.slides.add_slide(layout)
                slide.shapes.title.text = str(operation.get("title") or "")
                slide.placeholders[1].text = str(operation.get("body") or "")
        presentation.save(target)
        return "python-pptx"
    raise ValueError(f"Unsupported format: {suffix}")


def _apply_with_com(target: Path, operations: list[dict[str, Any]]) -> str:
    import pythoncom
    import win32com.client

    pythoncom.CoInitialize()
    try:
        suffix = target.suffix.lower()
        if suffix == ".docx":
            application = win32com.client.DispatchEx("Word.Application")
            application.Visible = False
            application.DisplayAlerts = 0
            document = application.Documents.Open(str(target))
            try:
                for operation in operations:
                    if operation["action"] == "replace_text":
                        find = document.Content.Find
                        find.Execute(
                            FindText=str(operation["find"]),
                            ReplaceWith=str(operation.get("replace") or ""),
                            Replace=2,
                        )
                    else:
                        document.Content.InsertAfter(
                            "\r" + str(operation.get("text") or "")
                        )
                document.Save()
            finally:
                document.Close(False)
                application.Quit()
            return "word_com"
        if suffix in {".xlsx", ".xlsm"}:
            application = win32com.client.DispatchEx("Excel.Application")
            application.Visible = False
            application.DisplayAlerts = False
            workbook = application.Workbooks.Open(str(target))
            try:
                for operation in operations:
                    cell = workbook.Worksheets(str(operation["sheet"])).Range(
                        str(operation["cell"])
                    )
                    if operation["action"] == "set_formula":
                        cell.Formula = str(operation.get("formula") or "")
                    else:
                        cell.Value = operation.get("value")
                workbook.RefreshAll()
                application.CalculateFull()
                workbook.Save()
            finally:
                workbook.Close(False)
                application.Quit()
            return "excel_com"
        if suffix == ".pptx":
            application = win32com.client.DispatchEx("PowerPoint.Application")
            presentation = application.Presentations.Open(
                str(target), WithWindow=False
            )
            try:
                for operation in operations:
                    if operation["action"] == "replace_text":
                        for slide in presentation.Slides:
                            for shape in slide.Shapes:
                                if getattr(shape, "HasTextFrame", False):
                                    text_range = shape.TextFrame.TextRange
                                    text_range.Text = text_range.Text.replace(
                                        str(operation["find"]),
                                        str(operation.get("replace") or ""),
                                    )
                    else:
                        slide = presentation.Slides.Add(
                            presentation.Slides.Count + 1, 2
                        )
                        slide.Shapes.Title.TextFrame.TextRange.Text = str(
                            operation.get("title") or ""
                        )
                        slide.Shapes.Placeholders(2).TextFrame.TextRange.Text = str(
                            operation.get("body") or ""
                        )
                presentation.Save()
            finally:
                presentation.Close()
                application.Quit()
            return "powerpoint_com"
        raise ValueError(f"Unsupported format: {suffix}")
    finally:
        pythoncom.CoUninitialize()


def handle_inspect(args: dict, **_: Any) -> str:
    started = time.monotonic()
    try:
        path = _resolve_document(str(args.get("path") or ""))
        result = _inspect(path)
        record_operational(
            OperationalEvent(
                event_type="document_inspect",
                capability="documents",
                success=True,
                duration_ms=int((time.monotonic() - started) * 1000),
                tool_name="document_inspect",
            )
        )
        return tool_result(result)
    except Exception as exc:
        return tool_error(str(exc))


def handle_extract(args: dict, **_: Any) -> str:
    started = time.monotonic()
    try:
        path = _resolve_document(str(args.get("path") or ""))
        maximum = max(1000, min(50000, int(args.get("max_chars") or 12000)))
        result = _extract(path, str(args.get("selector") or "") or None, maximum)
        record_operational(
            OperationalEvent(
                event_type="document_extract",
                capability="documents",
                success=True,
                duration_ms=int((time.monotonic() - started) * 1000),
                tool_name="document_extract",
            )
        )
        return tool_result(result)
    except Exception as exc:
        return tool_error(str(exc))


def handle_plan(args: dict, **_: Any) -> str:
    try:
        path = _resolve_document(str(args.get("path") or ""))
        output_raw = str(args.get("output_path") or "").strip()
        output_path = _resolve_document(output_raw, must_exist=False) if output_raw else None
        operations = args.get("operations")
        if not isinstance(operations, list) or not operations:
            return tool_error("operations must contain at least one change.")
        return tool_result(_plan_changes(path, output_path, operations))
    except Exception as exc:
        return tool_error(str(exc))


def handle_apply(args: dict, **_: Any) -> str:
    started = time.monotonic()
    try:
        plan_path = _plan_path(str(args.get("plan_id") or ""))
        record = json.loads(plan_path.read_text(encoding="utf-8"))
        if record.get("used"):
            raise ValueError("This plan has already been applied.")
        if time.time() > float(record.get("expires_at") or 0):
            raise ValueError("This plan has expired; create a fresh preview.")
        source = _resolve_document(record["path"])
        if hashlib.sha256(source.read_bytes()).hexdigest() != record["source_sha256"]:
            raise ValueError("The source changed after planning; create a fresh plan.")
        if not _ask_approval(record):
            record_compliance(
                ComplianceEvent(
                    event_type="office_change",
                    action="apply",
                    success=False,
                    resource=str(source),
                    approval="denied",
                    tool_name="office_apply_changes",
                )
            )
            return tool_error("User approval was denied or unavailable.")

        output = (
            _resolve_document(record["output_path"], must_exist=False)
            if record.get("output_path")
            else source
        )
        backup = None
        if output == source:
            backup_dir = source.parent / ".hermes-backups"
            backup_dir.mkdir(exist_ok=True)
            backup = backup_dir / (
                f"{source.stem}-{time.strftime('%Y%m%d-%H%M%S')}{source.suffix}"
            )
            shutil.copy2(source, backup)
        else:
            output.parent.mkdir(parents=True, exist_ok=True)
        staging = output.parent / (
            f".{output.stem}.hermes-staging-{uuid.uuid4().hex}{output.suffix}"
        )
        shutil.copy2(source, staging)
        try:
            engine = _apply_operations(staging, record["operations"])
            verification = _inspect(staging)
            os.replace(staging, output)
            verification["path"] = str(output)
        finally:
            staging.unlink(missing_ok=True)
        record["used"] = True
        record["applied_at"] = time.time()
        plan_path.write_text(json.dumps(record, indent=2), encoding="utf-8")
        duration = int((time.monotonic() - started) * 1000)
        record_operational(
            OperationalEvent(
                event_type="office_change",
                capability="office",
                success=True,
                duration_ms=duration,
                tool_name="office_apply_changes",
            )
        )
        record_compliance(
            ComplianceEvent(
                event_type="office_change",
                action="apply",
                success=True,
                resource=str(output),
                approval="approved",
                tool_name="office_apply_changes",
            )
        )
        return tool_result(
            {
                "applied": True,
                "output_path": str(output),
                "backup_path": str(backup) if backup else None,
                "engine": engine,
                "verification": verification,
            }
        )
    except Exception as exc:
        return tool_error(str(exc))


TOOL_DEFINITIONS = (
    ("document_inspect", DOCUMENT_INSPECT_SCHEMA, handle_inspect, "📄"),
    ("document_extract", DOCUMENT_EXTRACT_SCHEMA, handle_extract, "📑"),
    ("office_plan_changes", OFFICE_PLAN_SCHEMA, handle_plan, "📝"),
    ("office_apply_changes", OFFICE_APPLY_SCHEMA, handle_apply, "✅"),
)
